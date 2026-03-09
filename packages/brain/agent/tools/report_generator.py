"""
Report Generator — multi-format content synthesis engine.

Generates structured reports from research findings in multiple formats:
  - Markdown reports
  - PDF documents (via weasyprint)
  - Slide decks (via python-pptx)
  - Self-contained web pages (HTML + CSS)

Designed to be used as a tool in the agent's tool registry.
"""

from __future__ import annotations

import logging
import os
from dataclasses import dataclass
from typing import Any, Optional

logger = logging.getLogger("brain.agent.tools.report_generator")


@dataclass
class ReportConfig:
    """Configuration for report generation."""
    title: str = "Research Report"
    format: str = "markdown"  # markdown | pdf | slides | webpage
    output_dir: str = "/tmp/kestrel-outputs"
    include_toc: bool = True
    include_sources: bool = True
    style: str = "professional"  # professional | academic | minimal


class ReportGenerator:
    """Generates multi-format reports from structured content."""

    def __init__(self, config: Optional[ReportConfig] = None):
        self._config = config or ReportConfig()

    async def generate(
        self,
        content: str,
        config: Optional[ReportConfig] = None,
    ) -> dict[str, Any]:
        """Generate a report in the specified format.

        Returns:
            dict with "path" (output file path), "format", and "size_bytes"
        """
        cfg = config or self._config
        os.makedirs(cfg.output_dir, exist_ok=True)

        if cfg.format == "pdf":
            return await self._generate_pdf(content, cfg)
        elif cfg.format == "slides":
            return await self._generate_slides(content, cfg)
        elif cfg.format == "webpage":
            return await self._generate_webpage(content, cfg)
        else:
            return await self._generate_markdown(content, cfg)

    async def _generate_markdown(self, content: str, cfg: ReportConfig) -> dict[str, Any]:
        """Generate a markdown report."""
        path = os.path.join(cfg.output_dir, f"{_safe_filename(cfg.title)}.md")
        report = f"# {cfg.title}\n\n{content}"
        with open(path, "w", encoding="utf-8") as f:
            f.write(report)
        return {"path": path, "format": "markdown", "size_bytes": os.path.getsize(path)}

    async def _generate_pdf(self, content: str, cfg: ReportConfig) -> dict[str, Any]:
        """Generate a PDF report via weasyprint."""
        path = os.path.join(cfg.output_dir, f"{_safe_filename(cfg.title)}.pdf")
        html = _content_to_html(content, cfg.title, cfg.style)

        try:
            from weasyprint import HTML
            HTML(string=html).write_pdf(path)
            return {"path": path, "format": "pdf", "size_bytes": os.path.getsize(path)}
        except ImportError:
            logger.warning("weasyprint not installed, falling back to markdown")
            return await self._generate_markdown(content, cfg)

    async def _generate_slides(self, content: str, cfg: ReportConfig) -> dict[str, Any]:
        """Generate a slide deck via python-pptx."""
        path = os.path.join(cfg.output_dir, f"{_safe_filename(cfg.title)}.pptx")

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation()

            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = cfg.title

            # Content slides from sections
            sections = content.split("\n## ")
            for section in sections:
                if not section.strip():
                    continue
                lines = section.strip().split("\n", 1)
                heading = lines[0].strip().lstrip("# ")
                body = lines[1].strip() if len(lines) > 1 else ""

                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = heading
                if body and len(slide.placeholders) > 1:
                    slide.placeholders[1].text = body[:500]

            prs.save(path)
            return {"path": path, "format": "slides", "size_bytes": os.path.getsize(path)}
        except ImportError:
            logger.warning("python-pptx not installed, falling back to markdown")
            return await self._generate_markdown(content, cfg)

    async def _generate_webpage(self, content: str, cfg: ReportConfig) -> dict[str, Any]:
        """Generate a self-contained HTML page."""
        path = os.path.join(cfg.output_dir, f"{_safe_filename(cfg.title)}.html")
        html = _content_to_html(content, cfg.title, cfg.style)
        with open(path, "w", encoding="utf-8") as f:
            f.write(html)
        return {"path": path, "format": "webpage", "size_bytes": os.path.getsize(path)}


def _safe_filename(title: str) -> str:
    """Convert a title to a safe filename."""
    return "".join(c if c.isalnum() or c in "-_ " else "" for c in title).strip().replace(" ", "_")[:80]


def _content_to_html(content: str, title: str, style: str = "professional") -> str:
    """Convert markdown content to styled HTML."""
    styles = {
        "professional": "font-family: 'Segoe UI', system-ui; color: #1a1a2e; background: #fff;",
        "academic": "font-family: 'Times New Roman', serif; color: #333; background: #fefefe;",
        "minimal": "font-family: monospace; color: #000; background: #fff;",
    }

    body_lines = []
    for line in content.split("\n"):
        if line.startswith("### "):
            body_lines.append(f"<h3>{line[4:]}</h3>")
        elif line.startswith("## "):
            body_lines.append(f"<h2>{line[3:]}</h2>")
        elif line.startswith("# "):
            body_lines.append(f"<h1>{line[2:]}</h1>")
        elif line.startswith("- "):
            body_lines.append(f"<li>{line[2:]}</li>")
        elif line.startswith("**") and line.endswith("**"):
            body_lines.append(f"<strong>{line[2:-2]}</strong>")
        elif line.strip():
            body_lines.append(f"<p>{line}</p>")

    body_html = "\n".join(body_lines)
    css = styles.get(style, styles["professional"])

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        body {{ {css} max-width: 800px; margin: 0 auto; padding: 2rem; line-height: 1.6; }}
        h1 {{ color: #2563eb; }}
        h2 {{ color: #1e40af; border-bottom: 2px solid #e5e7eb; padding-bottom: 0.5rem; }}
        h3 {{ color: #3730a3; }}
        li {{ margin: 0.25rem 0; }}
        p {{ margin: 0.75rem 0; }}
    </style>
</head>
<body>
<h1>{title}</h1>
{body_html}
<footer style="margin-top:3rem; padding-top:1rem; border-top:1px solid #e5e7eb; color:#6b7280; font-size:0.875rem;">
    Generated by Kestrel Deep Research Engine
</footer>
</body>
</html>"""
