"""
Content generation subgraph — AIGC pipeline for slides, web pages, PDFs.

**STATUS: BETA — Not production-ready.**

Known limitations:
  - generate_outline() returns a hardcoded structural template; it does not
    use an LLM to derive an outline from the source text.
  - draft_content() distributes source text chunks across outline sections
    instead of using an LLM to write section-specific content.
  - review_content() checks only file existence and byte size; it does not
    perform any qualitative content review.

Do not expose this subgraph as a production feature until the placeholder
stages are replaced with real LLM-backed implementations.

Graph topology:
  START → outline → draft → format → review → END
"""

from __future__ import annotations

import logging
import os
from typing import Any

from agent.runtime.state import ContentState

logger = logging.getLogger("brain.agent.runtime.graphs.content")


async def generate_outline(state: ContentState) -> dict[str, Any]:
    """Generate a content outline from source text.

    BETA: Returns a hardcoded structural template.
    TODO: Use an LLM to derive a topic-specific outline from source_text.

    Creates a structured outline appropriate for the target format:
    - slides: title slide + 8-12 content slides + summary
    - webpage: header, sections, footer
    - pdf: title page, TOC, chapters, conclusion
    """
    logger.warning(
        "BETA: generate_outline() is using a hardcoded template. "
        "LLM-based outline generation is not yet implemented."
    )
    content_type = state.get("content_type", "slides")

    # Default outline structure by type
    if content_type == "slides":
        outline = [
            {"type": "title", "heading": "Presentation Title", "content": ""},
            {"type": "content", "heading": "Overview", "content": "Key points"},
            {"type": "content", "heading": "Key Findings", "content": "Details"},
            {"type": "content", "heading": "Analysis", "content": "Insights"},
            {"type": "content", "heading": "Conclusion", "content": "Summary"},
        ]
    elif content_type == "webpage":
        outline = [
            {"type": "header", "heading": "Page Title", "content": ""},
            {"type": "section", "heading": "Introduction", "content": ""},
            {"type": "section", "heading": "Main Content", "content": ""},
            {"type": "section", "heading": "Conclusion", "content": ""},
            {"type": "footer", "heading": "", "content": ""},
        ]
    else:  # pdf
        outline = [
            {"type": "title_page", "heading": "Document Title", "content": ""},
            {"type": "chapter", "heading": "Introduction", "content": ""},
            {"type": "chapter", "heading": "Findings", "content": ""},
            {"type": "chapter", "heading": "Conclusion", "content": ""},
        ]

    return {"outline": outline}


async def draft_content(state: ContentState) -> dict[str, Any]:
    """Generate draft content for each outline section.

    BETA: Distributes source text chunks across sections verbatim.
    TODO: Use an LLM to write section-specific content from the outline.
    """
    outline = state.get("outline", [])
    source_text = state.get("source_text", "")

    # Placeholder: chunk source text across sections
    # In production, replace with an LLM call per section
    sections = source_text.split("\n\n") if source_text else [""]
    draft_parts = []

    for i, section in enumerate(outline):
        content = sections[i] if i < len(sections) else ""
        draft_parts.append(f"## {section.get('heading', 'Section')}\n\n{content}")

    return {"draft": "\n\n".join(draft_parts)}


async def format_output(state: ContentState) -> dict[str, Any]:
    """Format the draft into the target output format."""
    content_type = state.get("content_type", "slides")
    draft = state.get("draft", "")
    output_dir = state.get("output_dir", "/tmp/kestrel-outputs")
    parent_task_id = state.get("parent_task_id", "unknown")

    os.makedirs(output_dir, exist_ok=True)

    output_path = ""

    if content_type == "slides":
        output_path = os.path.join(output_dir, f"{parent_task_id}_slides.pptx")
        try:
            from pptx import Presentation

            prs = Presentation()
            sections = draft.split("## ")
            for section in sections:
                if not section.strip():
                    continue
                lines = section.strip().split("\n", 1)
                title = lines[0].strip()
                body = lines[1].strip() if len(lines) > 1 else ""
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = title
                if body and slide.placeholders[1]:
                    slide.placeholders[1].text = body[:500]
            prs.save(output_path)
        except ImportError:
            logger.warning("python-pptx not installed, saving as markdown")
            output_path = output_path.replace(".pptx", ".md")
            with open(output_path, "w") as f:
                f.write(draft)

    elif content_type == "webpage":
        output_path = os.path.join(output_dir, f"{parent_task_id}_page.html")
        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Generated Report</title>
    <style>
        body {{ font-family: system-ui; max-width: 800px; margin: 0 auto; padding: 2rem; }}
        h2 {{ color: #2563eb; border-bottom: 2px solid #e5e7eb; padding-bottom: 0.5rem; }}
    </style>
</head>
<body>
{_markdown_to_html(draft)}
</body>
</html>"""
        with open(output_path, "w") as f:
            f.write(html)

    elif content_type == "pdf":
        output_path = os.path.join(output_dir, f"{parent_task_id}_report.pdf")
        try:
            from weasyprint import HTML
            html_content = f"<html><body>{_markdown_to_html(draft)}</body></html>"
            HTML(string=html_content).write_pdf(output_path)
        except ImportError:
            logger.warning("weasyprint not installed, saving as markdown")
            output_path = output_path.replace(".pdf", ".md")
            with open(output_path, "w") as f:
                f.write(draft)

    else:
        # Default: save as markdown
        output_path = os.path.join(output_dir, f"{parent_task_id}_output.md")
        with open(output_path, "w") as f:
            f.write(draft)

    return {"formatted_output": output_path}


async def review_content(state: ContentState) -> dict[str, Any]:
    """Review generated content for quality.

    BETA: Checks file existence and byte size only — no qualitative review.
    TODO: Use an LLM to review content accuracy, completeness, and tone.
    """
    output_path = state.get("formatted_output", "")
    content_type = state.get("content_type", "")

    feedback = f"Content generated successfully as {content_type}: {output_path}"

    if output_path and os.path.exists(output_path):
        size = os.path.getsize(output_path)
        feedback += f" (size: {size} bytes)"

    return {"review_feedback": feedback}


def _markdown_to_html(md: str) -> str:
    """Simple markdown to HTML conversion."""
    lines = md.split("\n")
    html_lines = []
    for line in lines:
        if line.startswith("## "):
            html_lines.append(f"<h2>{line[3:]}</h2>")
        elif line.startswith("### "):
            html_lines.append(f"<h3>{line[4:]}</h3>")
        elif line.startswith("# "):
            html_lines.append(f"<h1>{line[2:]}</h1>")
        elif line.startswith("- "):
            html_lines.append(f"<li>{line[2:]}</li>")
        elif line.strip():
            html_lines.append(f"<p>{line}</p>")
    return "\n".join(html_lines)


def build_content_graph(*, checkpointer=None):
    """Build the content generation subgraph."""
    from langgraph.graph import END, START, StateGraph

    graph = StateGraph(ContentState)

    graph.add_node("outline", generate_outline)
    graph.add_node("draft", draft_content)
    graph.add_node("format", format_output)
    graph.add_node("review", review_content)

    graph.add_edge(START, "outline")
    graph.add_edge("outline", "draft")
    graph.add_edge("draft", "format")
    graph.add_edge("format", "review")
    graph.add_edge("review", END)

    compile_kwargs = {}
    if checkpointer:
        compile_kwargs["checkpointer"] = checkpointer

    return graph.compile(**compile_kwargs)
